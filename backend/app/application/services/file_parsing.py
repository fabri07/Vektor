"""Shared file parsing helpers for chat uploads and ingestion.

This module centralizes:
  - filename sanitization
  - MIME detection with secure extension fallback
  - canonical parsed_summary_json generation
  - text extraction for supported document formats
"""

from __future__ import annotations

import csv
import io
import re
from pathlib import Path
from typing import TYPE_CHECKING, Any

import filetype

if TYPE_CHECKING:
    from decimal import Decimal

from app.domain.expense_categories import strip_accents
from app.domain.header_keys import fold_header, match_key
from app.observability.logger import get_logger

logger = get_logger(__name__)

# Límite de seguridad (anti-DOS). Archivos > 16MB se rechazan con error claro.
# El manejo de archivos más grandes (streaming / re-arquitectura) queda para más
# adelante; con este techo el JSONB de uploaded_files no se acerca al límite de Neon.
MAX_FILE_SIZE_BYTES = 16 * 1024 * 1024  # 16 MB
MAX_FILE_SIZE_LABEL = "16 MB"

SAFE_FILENAME_RE = re.compile(r"[^a-zA-Z0-9.\-_]")

SPREADSHEET_MIMES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "text/csv",
}
TEXT_MIMES = {
    "text/plain",
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
IMAGE_MIMES = {
    "image/jpeg",
    "image/png",
    "image/heif",
    "image/heic",
}
ALLOWED_MIMES = SPREADSHEET_MIMES | TEXT_MIMES | IMAGE_MIMES

EXTENSION_TO_MIME = {
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".heic": "image/heic",
}

SUPPORTED_TYPES_LABEL = "xlsx, csv, txt, pdf, docx, pptx, jpg, png, heic"

# TODOS los sets de keywords de este módulo se comparan contra encabezados ya
# pasados por `fold_header` (minúsculas, sin acentos ni ñ), así que se declaran
# SIEMPRE en su forma sin tilde. Escribir además la variante acentuada no agrega
# cobertura: es código muerto que no puede matchear nunca.
# Columnas que indican transacciones de venta (monto cobrado)
VENTA_COLS = {
    "precio_venta",
    "venta",
    "ventas",
    "ingreso",
    "monto",
    "importe",
    "total_venta",
    "total_cobrado",
    "cobro",
}
# "precio" y "total" solos son ambiguos (también aparecen en inventarios)
# — se pesan por separado en infer_spreadsheet_type

GASTO_COLS = {"costo", "gasto", "gastos", "egreso", "compra", "deuda", "pago", "proveedor"}

# Señales fuertes: inequívocamente un catálogo/inventario — no aparecen en transacciones
CATALOGO_COLS = {"sku", "codigo", "inventario", "articulo", "item"}

# Señales débiles: pueden aparecer en ventas/gastos también (descripcion de venta, nombre del
# proveedor)
NOMBRE_COLS = {"producto", "nombre"}

# PRODUCTO_COLS = unión para retrocompatibilidad con código que ya lo usa
PRODUCTO_COLS = CATALOGO_COLS | NOMBRE_COLS | {"stock", "descripcion"}
FECHA_COLS = {"fecha", "date", "dia", "mes", "periodo"}

# Fechas que NO son de una operación: son un atributo de la persona. Se restan de
# la señal de fecha TRANSACCIONAL (`has_fecha_transaccional`), que es la que usa
# la regla de maestros para descartar un documento con movimientos. Sin esto, un
# maestro de clientes con "Fecha de nacimiento" activaba `has_fecha` (match por
# substring de "fecha") y no podía clasificar como clientes — aunque
# `fecha_nacimiento` esté listada como señal de cliente más abajo.
FECHA_NO_TRANSACCIONAL_COLS = {"nacimiento", "cumpleanos", "cumple"}

# FASE 3: señales de compra de mercadería/insumos para reventa (→ inventario, no gasto).
# Conservador: solo se rerutea a stock si HAY mercadería Y una columna de cantidad.
MERCADERIA_COLS = {"mercaderia", "insumo", "insumos", "reposicion"}
CANTIDAD_COLS = {"cantidad", "unidades", "unidad", "qty", "cant"}

# FASE 3: clasificación CONTEXTUAL venta vs gasto. Las columnas de dinero genéricas
# (monto/importe/total/precio/valor) son NEUTRALES — aparecen en cualquier documento
# financiero y no deben favorecer ventas por sí solas. El tipo se decide por señales
# FUERTES de contexto (scoring), y ante empate/ausencia → "general" (el usuario confirma).
MONEY_COLS = {"monto", "importe", "total", "precio", "valor", "monto_total", "importe_total"}
# Señales fuertes de venta (cliente, ticket, factura emitida, cobro, caja, medio de pago).
# La FORMA DE PAGO queda deliberadamente afuera: un libro de compras también trae
# "forma_pago"/"medio_pago", así que no discrimina venta de gasto — solo dice que
# hay una operación. Cuando estaba acá le empataba el score a una compra real
# (proveedor=1 vs forma_pago=1) y la regla -1 no disparaba: el libro de compras
# se importaba como catálogo y se perdían el COGS y la salida de caja. La señal
# vive en FORMA_PAGO_COLS/`has_forma_pago`, que es lo que usa esa regla.
VENTA_SIGNAL_COLS = {
    "venta", "ventas", "vendido", "vendida", "ingreso", "ingresos", "facturacion",
    "factura_emitida", "ticket", "cliente", "consumidor", "cobro",
    "cobrado", "caja", "fecha_venta",
}
# Señales fuertes de gasto/egreso (proveedor, categoría, concepto, servicio, etc.).
# Nota: "pago" suelto NO se incluye — colisiona con "metodo/medio_pago" (señal de venta).
GASTO_SIGNAL_COLS = {
    "gasto", "gastos", "egreso", "egresos", "proveedor", "categoria",
    "rubro", "concepto", "servicio", "alquiler", "sueldo", "salario", "impuesto",
    "honorarios", "mantenimiento", "comision", "flete", "logistica",
    "factura_recibida", "compra", "costo", "deuda",
}

# Señales transaccionales para desambiguar un LIBRO DE COMPRAS de un CATÁLOGO.
# Un catálogo de productos NO trae monto de transacción + método de pago + proveedor
# + fecha juntos; un libro de compras de mercadería SÍ. Esto evita que `sku/cantidad/
# costo_unitario` clasifiquen una compra como "stock" perdiendo el gasto (COGS) y la
# salida de caja.
# Columna de MONTO de transacción (lo que se pagó/cobró en la operación). "precio"/
# "valor" quedan FUERA a propósito: en un catálogo "precio" es precio de lista, no un
# monto de operación.
TRANSACCION_MONTO_COLS = {"total", "monto", "importe", "monto_total", "importe_total"}
# Columna de método/forma de pago (efectivo, transferencia, tarjeta…). Se matchea
# contra la CLAVE del header (`match_key`), que colapsa preposiciones Y acentos,
# así que alcanza con la forma canónica sin tilde: "Medio de Pago" entra por
# "medio_pago" y "Método de Pago" por "metodo_pago".
FORMA_PAGO_COLS = {"forma_pago", "medio_pago", "metodo_pago"}
# Columna de proveedor (contraparte de una compra). Se matchea por SUBSTRING: en
# los archivos reales el header casi nunca viene pelado — "Proveedor (tal cual se
# anotó)", "Nombre del Proveedor" o "Razón Social Proveedor" son todos la misma
# columna. Con coincidencia exacta, un libro de compras con el header decorado
# perdía la señal de contraparte y sólo entraba en la regla -1 si además traía
# "Forma de Pago"; sin esa columna volvía a clasificarse como catálogo.
PROVEEDOR_COLS = {"proveedor", "proveedores"}
# Nº de comprobante de la operación (remito, factura, recibo). Es el discriminante
# más fuerte contra un catálogo: una lista de precios no numera comprobantes. Se
# usa como evidencia ALTERNATIVA al monto de la operación — un libro de compras
# que factura por remito puede no traer una columna "Total".
# "facturacion"/"factura_emitida" quedan afuera a propósito: son columnas de total
# facturado (señal de VENTA), no el identificador del documento.
COMPROBANTE_COLS = {"remito", "factura", "comprobante", "recibo"}
COMPROBANTE_EXCLUDE_KEYS = {"facturacion", "factura_emitida"}

# Columnas que nombran un PRECIO. Aparecen en el catálogo tanto como en una
# transacción, así que una señal de venta/gasto que salga de una de ellas
# ("precio_venta" matchea "venta") no cuenta como contexto de operación.
PRECIO_COL_KEYS = {"precio", "costo", "valor"}

# F7a: señales de identidad fiscal/contacto para maestros de CLIENTES/PROVEEDORES
# (aditivo — no implementa import todavía, solo detección/mapeo). Discriminadores
# específicos de cada entidad, alineados a los campos que persisten los modelos
# Customer (dni/cuit) y Supplier (cuil, sin dni/cuit): un maestro de clientes trae
# "dni"/"cliente"; uno de proveedores trae "cuil"/"proveedor". "cuit" queda fuera
# de ambos sets discriminadores porque lo usan tanto empresas-cliente como
# proveedores — no desempata por sí solo.
# "documento" entra acá porque es como identifica a la persona buena parte de los
# maestros reales (columna "Documento" + columna "Tipo"), no "DNI". Riesgo asumido:
# un maestro de PROVEEDORES que use "Documento" en vez de "CUIL" va a caer en
# clientes — corregible desde el selector de sección de la hoja, y mejor que el
# "stock" que daba antes.
CLIENTE_SIGNAL_COLS = {
    "cliente", "clientes", "consumidor", "dni", "documento", "cumpleanos",
    "fecha_nacimiento",
}
PROVEEDOR_MASTER_COLS = {"proveedor", "proveedores", "cuil", "contacto"}
# Contacto/fiscal genérico: no discrimina cliente vs proveedor por sí solo, pero
# junto con "nombre" y la ausencia de señales transaccionales indica que la hoja
# es un maestro de identidad (no una lista de precios ni un libro de compras).
IDENTIDAD_CONTACTO_COLS = {
    "cuit", "razon_social", "email", "telefono",
    "direccion", "localidad", "provincia", "codigo_postal", "cp",
    "apellido",
}

VENTA_CTX = {"venta", "ingreso", "cobro", "ticket", "recibo", "pago recibido", "cobrado"}
GASTO_CTX = {"gasto", "compra", "pago", "factura", "proveedor", "egreso", "gaste"}
STOCK_CTX = {"stock", "inventario", "unidades", "cantidad", "mercaderia"}

AMOUNT_RE = re.compile(r"\$\s*[\d.,]+")

# Strings que representan ausencia de valor en imports
_NULL_STRINGS = {"", "nan", "none", "null", "n/a", "na", "-", "nd"}

# Umbral de % de nulls por columna a partir del cual se advierte al usuario
NULL_COLUMN_WARN_THRESHOLD = 0.35


def normalize_numeric(
    value: object,
    *,
    required: bool = False,
    field_label: str = "campo",
) -> Decimal | None:
    """Normaliza un valor numérico de import o API.

    - None / string vacío / strings nulos → None (o 422 si required)
    - math.nan / math.inf → None (o 422 si required)
    - Strings con $, comas, puntos → parseo ARS (1.234,56 → 1234.56)
    """
    import math  # noqa: PLC0415
    from decimal import Decimal, InvalidOperation  # noqa: PLC0415

    if value is None:
        if required:
            raise ValueError(f"{field_label} es obligatorio y no puede ser nulo.")
        return None

    if isinstance(value, float) and (math.isnan(value) or math.isinf(value)):
        if required:
            raise ValueError(f"{field_label} contiene un valor inválido (NaN/Inf).")
        return None

    str_val = str(value).strip().lower()
    if str_val in _NULL_STRINGS:
        if required:
            raise ValueError(f"{field_label} es obligatorio.")
        return None

    # Normalización de formato ARS: "1.234,56" → "1234.56"
    if "," in str_val and "." in str_val:
        str_val = str_val.replace(".", "").replace(",", ".")
    elif "," in str_val:
        str_val = str_val.replace(",", ".")
    str_val = str_val.lstrip("$").strip()

    try:
        return Decimal(str_val)
    except InvalidOperation as exc:
        if required:
            raise ValueError(
                f"{field_label} tiene un formato numérico inválido: {value!r}"
            ) from exc
        return None


def normalize_categorical(
    value: object,
    *,
    required: bool = False,
    default: str | None = None,
    field_label: str = "campo",
) -> str | None:
    """Normaliza un valor categórico de import o API."""
    if value is None:
        return (
            default
            if not required
            else (_ for _ in ()).throw(ValueError(f"{field_label} es obligatorio."))
        )
    str_val = str(value).strip()
    if str_val.lower() in _NULL_STRINGS:
        if required:
            raise ValueError(f"{field_label} es obligatorio.")
        return default
    return str_val or (
        default
        if not required
        else (_ for _ in ()).throw(ValueError(f"{field_label} no puede estar vacío."))
    )


def compute_column_null_stats(rows: list[dict[str, Any]]) -> dict[str, float]:
    """Calcula el porcentaje de valores nulos por columna en una lista de dicts.

    Retorna {col: null_pct} para cada columna presente en al menos una fila.
    """
    if not rows:
        return {}
    all_keys: set[str] = set()
    for row in rows:
        all_keys.update(row.keys())

    stats: dict[str, float] = {}
    total = len(rows)
    for col in all_keys:
        null_count = sum(
            1
            for row in rows
            if row.get(col) is None or str(row.get(col, "")).strip().lower() in _NULL_STRINGS
        )
        stats[col] = null_count / total
    return stats


def flag_columns_at_risk(
    null_stats: dict[str, float],
    threshold: float = NULL_COLUMN_WARN_THRESHOLD,
) -> list[dict[str, Any]]:
    """Retorna lista de columnas que superan el umbral de nulls, con recomendación."""
    return [
        {"column": col, "null_pct": round(pct, 4), "recommendation": "drop"}
        for col, pct in null_stats.items()
        if pct > threshold
    ]


def impute_column(values: list[Any], field_type: str) -> list[Any]:
    """Imputa valores nulos en una columna según su tipo.

    - field_type='quantity': nulos → 0
    - field_type='numeric':  nulos → mediana de los valores válidos
    - field_type='categorical': nulos → None (sin imputación)
    """
    import math  # noqa: PLC0415
    from decimal import Decimal  # noqa: PLC0415

    def _is_null(v: object) -> bool:
        if v is None:
            return True
        if isinstance(v, float) and (math.isnan(v) or math.isinf(v)):
            return True
        return str(v).strip().lower() in _NULL_STRINGS

    if field_type == "quantity":
        return [0 if _is_null(v) else v for v in values]

    if field_type == "numeric":
        median = impute_column_median(values)
        fill = median if median is not None else Decimal("0")
        return [fill if _is_null(v) else v for v in values]

    # categorical — sin imputación
    return [None if _is_null(v) else v for v in values]


def impute_column_median(values: list[Any]) -> Decimal | None:
    """Calcula la mediana de una lista de valores numéricos.
    Usa mediana (resistente a outliers de precios) en vez de media.
    Retorna None si no hay valores válidos.
    """
    import math  # noqa: PLC0415
    from decimal import Decimal  # noqa: PLC0415

    nums: list[Decimal] = []
    for v in values:
        if v is None:
            continue
        if isinstance(v, float) and (math.isnan(v) or math.isinf(v)):
            continue
        try:
            nums.append(Decimal(str(v)))
        except Exception:
            continue

    if not nums:
        return None
    nums_sorted = sorted(nums)
    mid = len(nums_sorted) // 2
    if len(nums_sorted) % 2 == 0:
        return (nums_sorted[mid - 1] + nums_sorted[mid]) / 2
    return nums_sorted[mid]


def sanitize_filename(filename: str) -> str:
    """Remove path traversal and unsafe characters from a filename."""
    filename = filename.replace("\\", "/").split("/")[-1]
    filename = SAFE_FILENAME_RE.sub("_", filename)
    return filename or "upload"


def detect_supported_mime(content: bytes, filename: str) -> str:
    """Detect a supported MIME type from content and filename."""
    kind = filetype.guess(content[:2048])
    ext = Path(filename).suffix.lower()

    detected = kind.mime if kind is not None else ""
    if detected == "application/zip" and ext in EXTENSION_TO_MIME:
        detected = EXTENSION_TO_MIME[ext]

    if not detected and ext in EXTENSION_TO_MIME:
        detected = EXTENSION_TO_MIME[ext]

    if detected == "image/jpg":
        detected = "image/jpeg"

    if detected not in ALLOWED_MIMES:
        raise ValueError(
            f"Tipo de archivo no soportado: {detected or 'desconocido'}. "
            f"Tipos aceptados: {SUPPORTED_TYPES_LABEL}."
        )
    return detected


def infer_source_format(filename: str, mime: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext:
        return ext.lstrip(".")
    if mime == "text/plain":
        return "txt"
    if mime == "text/csv":
        return "csv"
    return mime.rsplit("/", 1)[-1]


def analyze_headers(headers: list[str]) -> dict[str, Any]:
    """Classify headers and infer spreadsheet type."""
    # `fold_header` saca acentos y ñ además de bajar a minúsculas: los sets de
    # keywords de abajo se declaran SIN tilde y matchean igual una hoja que
    # escriba "Descripción", "Mercadería" o "Cumpleaños".
    normalized = [fold_header(h) for h in headers]
    # Clave heurística (sin preposiciones) para las señales que matchean EXACTO:
    # "medio_de_pago" y "medio_pago" son el mismo header. Se usa solo donde hace
    # falta — el resto de las señales sigue matcheando contra `normalized`.
    collapsed = [match_key(col) for col in normalized]

    has_fecha = any(any(k in col for k in FECHA_COLS) for col in normalized)
    # Fecha de OPERACIÓN: la de nacimiento/cumpleaños es un atributo de la persona
    # y no evidencia de que el documento tenga movimientos. Señal separada a
    # propósito: `has_fecha` conserva su significado para el resto del sistema.
    has_fecha_transaccional = any(
        any(k in col for k in FECHA_COLS)
        and not any(k in col for k in FECHA_NO_TRANSACCIONAL_COLS)
        for col in normalized
    )
    # FASE 3: venta/gasto por señales FUERTES de contexto (no por columna de dinero).
    venta_score = sum(any(k in col for k in VENTA_SIGNAL_COLS) for col in normalized)
    gasto_score = sum(any(k in col for k in GASTO_SIGNAL_COLS) for col in normalized)
    # Contexto de OPERACIÓN: una señal de venta/gasto que NO venga de una columna de
    # precio. "precio_venta" matchea "venta" por substring, pero es cómo se llama un
    # campo del catálogo — no evidencia de que haya una venta registrada. Sin este
    # filtro, una lista de precios con fecha de alta y una columna "Total" (valuación
    # del stock) tendría las tres señales de operación y se importaría como ventas,
    # inventando facturación.
    has_contexto_operacion = any(
        (any(k in col for k in VENTA_SIGNAL_COLS) or any(k in col for k in GASTO_SIGNAL_COLS))
        and not any(p in col for p in PRECIO_COL_KEYS)
        for col in normalized
    )
    has_venta = venta_score > 0
    has_gasto = gasto_score > 0
    has_producto = any(any(k in col for k in PRODUCTO_COLS) for col in normalized)
    # Señal fuerte de catálogo: sku/codigo/inventario/articulo/item — inequívocamente no transacción
    has_catalogo_fuerte = any(any(k in col for k in CATALOGO_COLS) for col in normalized)
    # Señal de nombre: producto/nombre — puede aparecer en ventas/gastos también
    has_nombre = any(any(k in col for k in NOMBRE_COLS) for col in normalized)
    # FASE 3: señales de compra de mercadería + cantidad (inventario para reventa)
    has_mercaderia = any(any(k in col for k in MERCADERIA_COLS) for col in normalized)
    has_cantidad = any(any(k in col for k in CANTIDAD_COLS) for col in normalized)

    # Señales ambiguas: "precio" / "total" solos pueden ser precio de catálogo
    has_precio_ambiguo = any(col in ("precio", "total", "price", "valor") for col in normalized)

    # Señales transaccionales de COMPRA (libro de compras ≠ catálogo): monto de la
    # operación + método de pago + proveedor. Se exige coincidencia exacta de columna
    # para no disparar con sufijos espurios.
    has_monto_transaccion = any(col in TRANSACCION_MONTO_COLS for col in normalized)
    has_forma_pago = any(key in FORMA_PAGO_COLS for key in collapsed)
    has_proveedor = any(any(k in col for k in PROVEEDOR_COLS) for col in normalized)
    has_comprobante = any(
        any(k in col for k in COMPROBANTE_COLS)
        and not any(x in col for x in COMPROBANTE_EXCLUDE_KEYS)
        for col in normalized
    )

    # F7a: señales de maestro de clientes/proveedores (identidad fiscal/contacto).
    cliente_score = sum(any(k in col for k in CLIENTE_SIGNAL_COLS) for col in normalized)
    proveedor_master_score = sum(
        any(k in col for k in PROVEEDOR_MASTER_COLS) for col in normalized
    )
    has_identidad_contacto = any(
        any(k in col for k in IDENTIDAD_CONTACTO_COLS) for col in normalized
    )

    inferred_type = infer_spreadsheet_type(
        has_fecha=has_fecha,
        has_fecha_transaccional=has_fecha_transaccional,
        has_contexto_operacion=has_contexto_operacion,
        has_venta=has_venta,
        has_gasto=has_gasto,
        has_producto=has_producto,
        has_precio_ambiguo=has_precio_ambiguo,
        has_catalogo_fuerte=has_catalogo_fuerte,
        has_nombre=has_nombre,
        has_mercaderia=has_mercaderia,
        has_cantidad=has_cantidad,
        venta_score=venta_score,
        gasto_score=gasto_score,
        has_monto_transaccion=has_monto_transaccion,
        has_forma_pago=has_forma_pago,
        has_proveedor=has_proveedor,
        has_comprobante=has_comprobante,
        cliente_score=cliente_score,
        proveedor_master_score=proveedor_master_score,
        has_identidad_contacto=has_identidad_contacto,
    )

    has_catalogo_signal = has_catalogo_fuerte or has_nombre
    confidence = "HIGH" if (has_fecha and has_venta and not has_catalogo_signal) else "MEDIUM"

    return {
        "has_fecha": has_fecha,
        "has_venta": has_venta,
        "has_gasto": has_gasto,
        "has_producto": has_producto,
        "inferred_type": inferred_type,
        "confidence": confidence,
    }


def infer_spreadsheet_type(
    *,
    has_fecha: bool,
    has_venta: bool,
    has_gasto: bool,
    has_producto: bool,
    has_precio_ambiguo: bool = False,
    has_catalogo_fuerte: bool = False,
    has_nombre: bool = False,
    has_mercaderia: bool = False,
    has_cantidad: bool = False,
    venta_score: int = 0,
    gasto_score: int = 0,
    has_monto_transaccion: bool = False,
    has_forma_pago: bool = False,
    has_proveedor: bool = False,
    has_comprobante: bool = False,
    cliente_score: int = 0,
    proveedor_master_score: int = 0,
    has_identidad_contacto: bool = False,
    has_fecha_transaccional: bool | None = None,
    has_contexto_operacion: bool | None = None,
) -> str:
    """Determina el tipo más probable del archivo tabular.

    Reglas (en orden de prioridad):
    -1. LIBRO DE COMPRAS: catálogo (sku/producto/cantidad) + (monto de transacción O
        nº de comprobante) + fecha + (forma de pago O proveedor) Y contexto de gasto
        dominante → gastos, aunque traiga sku/cantidad/costo_unitario. Una compra de
        mercadería es a la vez gasto (COGS) y salida de caja; el catálogo no la captura.
        CONSERVADOR: solo se dispara cuando, sin esta regla, el archivo caería en
        "stock" por error.
    -0.5 (F7a). MAESTRO DE CLIENTES/PROVEEDORES: señal de identidad fiscal/contacto
        (dni/documento/cliente para clientes; cuil/proveedor para proveedores) SIN
        ninguna señal transaccional (monto de operación, cantidad, fecha de
        OPERACIÓN — una fecha de nacimiento no cuenta) NI de catálogo fuerte
        (sku/codigo/inventario/articulo/item — regla 1) → clientes/proveedores.
        Disjunta con la regla -1 (esa exige monto+fecha; esta los excluye), así el
        orden entre ambas no importa. El guard de catálogo evita que un CATÁLOGO de
        productos con una columna "proveedor"/"cliente" (ej. quién lo distribuye) se
        confunda con un maestro — sigue siendo "stock" vía la regla 1. Corre ANTES
        que las reglas de catálogo débiles (2-4, basadas en "nombre") para que un
        maestro con columna "nombre" no se confunda con una lista de precios.
    0. Compra de mercadería/insumos + cantidad → inventario (FASE 3, conservador).
    1. Señal fuerte de catálogo (sku/codigo/inventario/articulo/item) → stock, salvo
       que haya evidencia transaccional completa (misma excepción que la regla 5:
       "Artículo"/"Item" también nombran el ítem vendido en un libro de ventas).
    2. Señal de nombre/producto sin venta explícita → stock.
    3. Señal de nombre/producto sin fecha → stock (lista de precios, catálogo).
    4. Señal de nombre/producto + precio ambiguo (no venta transaccional) → stock.
    5. Señal de nombre/producto CON evidencia transaccional completa (fecha + monto
       de la operación + contexto de venta/gasto) → sigue al scoring, no a stock.
       Sin la evidencia completa → stock.
    6. Sin señales de catálogo: fecha + venta → ventas; fecha + gasto → gastos.
    7. Fallbacks por señales sueltas.
    """
    # Los callers que no distinguen la fecha de operación de la de nacimiento, ni el
    # contexto fuerte del que sale de una columna de precio, mantienen el
    # comportamiento previo.
    if has_fecha_transaccional is None:
        has_fecha_transaccional = has_fecha
    if has_contexto_operacion is None:
        has_contexto_operacion = has_venta or has_gasto
    # Regla -1 (CONSERVADORA): un LIBRO DE COMPRAS de mercadería tiene a la vez columnas
    # de catálogo (sku/producto/cantidad) Y de operación (monto de transacción + fecha +
    # forma de pago/proveedor). Sin esta regla, las señales de catálogo lo clasificarían
    # como "stock" y se perdería el gasto (COGS) y la salida de caja. Solo se dispara si:
    #   (a) hay señal de catálogo (sku/articulo/… o producto/nombre o cantidad) — si no la
    #       hay, el scoring venta/gasto de abajo ya resuelve bien y no hay que intervenir;
    #   (b) hay evidencia de operación documentada (monto de transacción O nº de
    #       comprobante) + fecha + (forma de pago O proveedor) — patrón de operación, no
    #       de lista de precios; y
    #   (c) el contexto de GASTO domina al de VENTA (gasto_score > venta_score) — así una
    #       venta con cliente+forma_pago+total NO cae acá.
    # Un catálogo real (precio_venta/stock_actual/stock_minimo, sin forma_pago/proveedor
    # ni monto de operación) NO cumple (b) y sigue siendo "stock".
    #
    # El comprobante entra como alternativa al monto porque un libro de compras que
    # factura por remito puede no traer una columna "Total" — y numerar remitos es algo
    # que un catálogo no hace nunca. Sin esa alternativa, la regla dependía de que el
    # archivo trajera JUSTO las columnas pelada "total" y "forma de pago": un libro real
    # (fecha + nº de remito + proveedor + cantidad + costo) volvía a caer en "stock",
    # perdiendo el COGS, la salida de caja y los movimientos de stock fechados que
    # respaldan las ventas del mismo archivo.
    has_catalogo_signal_any = has_catalogo_fuerte or has_nombre or has_cantidad
    has_operacion_documentada = has_monto_transaccion or has_comprobante
    if (
        has_catalogo_signal_any
        and has_operacion_documentada
        and has_fecha
        and (has_forma_pago or has_proveedor)
        and gasto_score > venta_score
    ):
        return "gastos"

    # Regla -0.5 (F7a, CONSERVADORA): un maestro de CLIENTES o PROVEEDORES trae
    # identidad fiscal/contacto (dni/cliente, cuil/proveedor, +opcionalmente
    # nombre+email/telefono/localidad/etc.) pero NINGUNA señal transaccional
    # (monto de operación, cantidad, fecha) NI de catálogo (sku/codigo/inventario/
    # articulo/item — regla 1, `has_catalogo_fuerte`). Sin esta regla, un maestro
    # con columna "nombre" caería en "stock" (regla 2-4) o, si trae "proveedor", en
    # "gastos" (GASTO_SIGNAL_COLS) — perdiendo el maestro. El guard de catálogo
    # evita el caso inverso: un CATÁLOGO de productos que además trae una columna
    # "proveedor"/"cliente" (p.ej. quién lo distribuye) NO es un maestro de
    # clientes/proveedores — sigue siendo "stock" (regla 1, más abajo). Un maestro
    # real no trae sku/codigo/inventario/articulo/item. Empate o ausencia de
    # discriminador → no se adivina (cae a las reglas siguientes, comportamiento
    # previo).
    has_maestro_signal = (
        cliente_score > 0
        or proveedor_master_score > 0
        or (has_nombre and has_identidad_contacto)
    )
    if (
        has_maestro_signal
        and not has_monto_transaccion
        and not has_cantidad
        and not has_fecha_transaccional
        and not has_catalogo_fuerte
    ):
        if proveedor_master_score > cliente_score:
            return "proveedores"
        if cliente_score > proveedor_master_score:
            return "clientes"

    # FASE 3: compra de mercadería/insumos para reventa CON columna de cantidad →
    # inventario, NO gasto corriente. Conservador: requiere AMBAS señales para no
    # absorber "compra de servicios/alquiler" (sin columnas de inventario).
    if has_mercaderia and has_cantidad:
        return "stock"

    # Evidencia de que el archivo registra OPERACIONES y no describe un catálogo:
    # fecha + monto DE LA TRANSACCIÓN + contexto de venta/gasto que no salga de una
    # columna de precio. Las tres juntas, y ninguna alcanza sola: un catálogo tiene
    # fecha de alta, y "precio_venta" activa el contexto por substring.
    #
    # La usan las reglas 1 y 5 con el mismo significado, y por eso se calcula una
    # sola vez: son la misma pregunta hecha sobre dos señales de catálogo distintas
    # (la fuerte y la de nombre), y tenerlas escritas por separado fue justamente
    # lo que las dejó divergir.
    evidencia_transaccional = has_fecha and has_monto_transaccion and has_contexto_operacion

    # Señal fuerte (sku, inventario, articulo, codigo, item) → catálogo, SALVO que
    # el archivo traiga la evidencia transaccional completa.
    #
    # La excepción existe porque la regla no era cierta como estaba escrita
    # ("inequívocamente catálogo"): `articulo` e `item` son cómo se llama la columna
    # del ítem VENDIDO en media exportación de ventas de un kiosco, no sólo la del
    # catálogo. Un libro de ventas con "Artículo" volvía stock y no se importaba
    # ninguna venta, mientras que el MISMO archivo con la columna llamada "Producto"
    # entraba bien — porque la señal de nombre (regla 5) sí tenía esta excepción y
    # la fuerte no. La asimetría estaba tapada por los acentos: "Artículo" con tilde
    # no matcheaba el keyword `articulo` y se salvaba de rebote.
    if has_catalogo_fuerte and not evidencia_transaccional:
        return "stock"

    # Nombre/producto sin venta explícita → catálogo (descripcion en ventas suele ir con monto)
    if has_nombre and not has_venta:
        return "stock"

    # Nombre/producto sin fecha → lista de precios/catálogo
    if has_nombre and not has_fecha:
        return "stock"

    # Nombre/producto + precio ambiguo (ej: nombre+precio sin monto/importe) → catálogo
    if has_nombre and has_precio_ambiguo and not has_venta:
        return "stock"

    # Nombre/producto con las tres señales de OPERACIÓN juntas (fecha + monto de la
    # transacción + contexto de venta/gasto que no salga de una columna de precio) no
    # es ambiguo: es un libro de ventas o de gastos con columna de producto, y sigue al
    # scoring de abajo. Sin alguna de las tres, la señal de nombre gana y es catálogo.
    #
    # Este era el desempate "ante la duda, stock". Pero las tres reglas de arriba ya
    # filtran todo lo demás, así que este return solo era alcanzable con
    # nombre+venta+fecha — o sea que no desempataba un caso dudoso: pisaba exactamente
    # las exportaciones de ventas.
    #
    # El contexto tiene que ser fuerte: un catálogo con "precio_venta" ya activa
    # `has_venta` por substring, y si además trae fecha de alta y una columna "Total"
    # (valuación del stock) tendría las tres señales sin ser una transacción.
    if has_nombre and not evidencia_transaccional:
        return "stock"

    # FASE 3: discriminación venta vs gasto por CONTEXTO (scoring de señales fuertes).
    # Las columnas de dinero genéricas (monto/importe/total) NO cuentan acá. Empate de
    # señales o ausencia total → "general" (ambiguo): el usuario confirma el tipo, no se
    # importa como venta en silencio.
    if has_venta and has_gasto:
        if gasto_score > venta_score:
            return "gastos"
        if venta_score > gasto_score:
            return "ventas"
        return "general"  # señales de ambos, empate → ambiguo
    if has_gasto:
        return "gastos"
    if has_venta:
        return "ventas"
    # Sin señales fuertes de contexto (solo dinero/fecha/descripción) → ambiguo.
    return "general"


def rows_to_dicts(headers: list[str], rows: list[list[Any]]) -> list[dict[str, Any]]:
    """Mapea filas a dicts por header, tolerante a filas irregulares.

    Si una fila tiene menos celdas que headers, las faltantes quedan None
    (no se descartan en silencio como hacía `zip(strict=False)`). Filas None y
    celdas extra (más allá de los headers) se ignoran.
    """
    out: list[dict[str, Any]] = []
    n = len(headers)
    for row in rows:
        cells = list(row) if row is not None else []
        out.append(
            {
                headers[i]: (str(cells[i]) if i < len(cells) and cells[i] is not None else None)
                for i in range(n)
            }
        )
    return out


def _detect_header_row(rows: list[list[Any]], *, max_scan: int = 15) -> int:
    """Detecta el índice de la fila de encabezado, saltando títulos/filas vacías.

    Real-world: planillas exportadas suelen tener un título ("Ventas Junio") y/o
    filas en blanco arriba del encabezado real. Heurística: la primera fila con
    ≥2 celdas no vacías, mayoría de texto (no números), y cuya fila siguiente
    tenga una cantidad de celdas no vacías comparable (datos). Si ninguna
    califica, devuelve 0 (comportamiento previo).
    """

    def _nonempty(row: list[Any]) -> list[Any]:
        return [c for c in (row or []) if c is not None and str(c).strip() != ""]

    def _is_number(v: Any) -> bool:
        s = str(v).strip().replace(".", "").replace(",", "").replace("$", "").replace("%", "")
        return s.isdigit()

    limit = min(len(rows), max_scan)
    for i in range(limit):
        non_empty = _nonempty(rows[i])
        if len(non_empty) < 2:
            continue  # título o fila casi vacía
        numeric_ratio = sum(1 for c in non_empty if _is_number(c)) / len(non_empty)
        if numeric_ratio > 0.5:
            continue  # parece una fila de datos, no encabezados
        # La fila siguiente debería tener datos (≥1 celda no vacía).
        if i + 1 < len(rows) and len(_nonempty(rows[i + 1])) >= 1:
            return i
        if i + 1 >= len(rows):
            return i
    return 0


# ── Libro Diario (doble encabezado: Dinero/Mercadería × Ingreso/Egreso) ───────
# Formato típico de libros contables manuales en Excel: una fila superior con
# grupos en celdas combinadas ("Dinero", "Mercadería") y una segunda fila con
# "Ingreso"/"Egreso" por grupo. Muchas filas continúan la fecha de la anterior
# (forward-fill). Semántica contable (regla canónica del negocio):
#   - Dinero/Ingreso                 → venta (entra plata; la mercadería sale)
#   - Mercadería/Ingreso             → compra de mercadería (entra stock, sale
#                                      plata) → gasto INVENTORY (COGS)
#   - Dinero/Egreso sin mercadería   → gasto operativo (OPEX)
#   - Plata que entra Y sale en la misma fila, o salida de mercadería sin plata
#     → ambiguo: revisión manual ("otros"), nunca se adivina.

_LD_GRUPO_DINERO = {"dinero", "caja", "efectivo", "plata"}
_LD_GRUPO_MERCADERIA = {"mercaderia", "mercaderias", "stock", "inventario"}
_LD_SUB_INGRESO = {"ingreso", "ingresos", "entrada", "entradas"}
_LD_SUB_EGRESO = {"egreso", "egresos", "salida", "salidas"}
_LD_DETALLE_KEYS = {"detalle", "descripcion", "concepto", "glosa", "movimiento", "observa"}
# Filas de cierre que no son movimientos (subtotales del libro).
_LD_TOTAL_PREFIXES = ("total", "subtotal", "saldo", "acumulado")
# Hojas DERIVADAS: agregados que Véktor recalcula solo desde los movimientos
# (resúmenes por medio de pago, ganancias, balances). Importarlas además de las
# hojas de movimientos del mismo archivo suma esos totales por segunda vez —
# facturación fantasma, y encima con la fila "TOTAL" sumando otra vez las de
# arriba. Mismo criterio que el margen en `column_mapping_service`: lo que el
# sistema CALCULA no se importa como dato. No se descartan: se preservan sin
# clasificar y destildadas, así el usuario puede asignarles sección a mano si de
# verdad las quiere (la regla cambia el default, no el permiso).
_DERIVED_SHEET_PREFIXES = ("ganancia", "resumen", "balance")


def _ld_norm_cell(value: Any) -> str:
    """lower, sin tildes, separadores colapsados a un espacio."""
    if value is None:
        return ""
    s = strip_accents(str(value).strip().lower())
    return re.sub(r"[\s\-_/]+", " ", s)


def detect_libro_diario_header(
    rows: list[list[Any]], *, max_scan: int = 15
) -> tuple[int, dict[str, int]] | None:
    """Detecta el doble encabezado del Libro Diario.

    Busca una fila con un grupo "Dinero" (celdas combinadas → None en las
    columnas no-ancla) seguida de una fila con "Ingreso"/"Egreso". Devuelve
    ``(header_idx, col_map)`` con índices de columna para ``dinero_ingreso``,
    ``dinero_egreso``, ``mercaderia_ingreso``, ``mercaderia_egreso``, ``fecha``
    y ``detalle`` — o ``None`` si la estructura no está presente. Exige ambas
    columnas de Dinero para evitar falsos positivos.
    """
    limit = min(len(rows), max_scan)
    for i in range(limit - 1):
        top = [_ld_norm_cell(c) for c in (rows[i] or [])]
        if not any(c in _LD_GRUPO_DINERO for c in top):
            continue
        sub = [_ld_norm_cell(c) for c in (rows[i + 1] or [])]
        # Forward-fill de los grupos: las celdas combinadas llegan vacías.
        groups: list[str] = []
        current = ""
        for c in top:
            if c:
                current = c
            groups.append(current)
        col_map: dict[str, int] = {}
        for j in range(max(len(groups), len(sub))):
            g = groups[j] if j < len(groups) else ""
            s = sub[j] if j < len(sub) else ""
            if s in _LD_SUB_INGRESO or s in _LD_SUB_EGRESO:
                sfx = "ingreso" if s in _LD_SUB_INGRESO else "egreso"
                if g in _LD_GRUPO_DINERO:
                    col_map.setdefault(f"dinero_{sfx}", j)
                elif g in _LD_GRUPO_MERCADERIA:
                    col_map.setdefault(f"mercaderia_{sfx}", j)
                continue
            label = s or (top[j] if j < len(top) else "")
            if any(k in label for k in FECHA_COLS):
                col_map.setdefault("fecha", j)
            elif any(k in label for k in _LD_DETALLE_KEYS):
                col_map.setdefault("detalle", j)
        if "dinero_ingreso" in col_map and "dinero_egreso" in col_map:
            return i, col_map
    return None


def _ld_cell(cells: list[Any], col_map: dict[str, int], key: str) -> Any:
    idx = col_map.get(key)
    return cells[idx] if idx is not None and idx < len(cells) else None


def _ld_amount(value: Any) -> Decimal | None:
    """Monto positivo de la celda o None (vacío / no numérico / ≤ 0)."""
    amount = normalize_numeric(value)
    return amount if amount is not None and amount > 0 else None


def parse_libro_diario(
    rows: list[list[Any]], hdr_idx: int, col_map: dict[str, int]
) -> dict[str, list[dict[str, Any]]]:
    """Clasifica las filas del Libro Diario en ventas / gastos / otros.

    Emite dicts con claves canónicas que el import resuelve por keyword
    (``fecha``, ``detalle``, ``monto``, ``categoria``, ``forma_pago``) en vez
    de los encabezados compuestos crudos. La fecha se forward-fillea (las filas
    del libro continúan la fecha de la anterior).
    """
    ventas: list[dict[str, Any]] = []
    gastos: list[dict[str, Any]] = []
    otros: list[dict[str, Any]] = []
    last_fecha: Any = None

    for raw_row in rows[hdr_idx + 2 :]:
        cells = list(raw_row) if raw_row is not None else []
        fecha_raw = _ld_cell(cells, col_map, "fecha")
        if fecha_raw is not None and str(fecha_raw).strip() != "":
            last_fecha = fecha_raw
        detalle_raw = _ld_cell(cells, col_map, "detalle")
        detalle = (
            str(detalle_raw).strip()
            if detalle_raw is not None and str(detalle_raw).strip()
            else None
        )
        if detalle and _ld_norm_cell(detalle).startswith(_LD_TOTAL_PREFIXES):
            continue  # subtotales/saldos del libro: no son movimientos

        din_in = _ld_amount(_ld_cell(cells, col_map, "dinero_ingreso"))
        din_out = _ld_amount(_ld_cell(cells, col_map, "dinero_egreso"))
        mer_in = _ld_amount(_ld_cell(cells, col_map, "mercaderia_ingreso"))
        mer_out = _ld_amount(_ld_cell(cells, col_map, "mercaderia_egreso"))
        if not any((din_in, din_out, mer_in, mer_out)):
            continue

        fecha = str(last_fecha).strip() if last_fecha is not None else None
        base: dict[str, Any] = {"fecha": fecha, "detalle": detalle}

        if din_in and (din_out or mer_in):
            # Plata que entra y sale (o venta + compra) en la misma fila →
            # nunca adivinar: revisión manual.
            otros.append(
                {
                    **base,
                    "dinero_ingreso": str(din_in),
                    "dinero_egreso": str(din_out) if din_out else None,
                    "mercaderia_ingreso": str(mer_in) if mer_in else None,
                    "mercaderia_egreso": str(mer_out) if mer_out else None,
                }
            )
        elif din_in:
            ventas.append({**base, "monto": str(din_in)})
        elif mer_in:
            # Compra de mercadería: entra stock, sale plata → gasto INVENTORY
            # (COGS). Sin salida de plata en la fila, la compra fue a cuenta
            # corriente del proveedor (fiado).
            gasto = {**base, "monto": str(din_out or mer_in), "categoria": "Mercadería"}
            if not din_out:
                gasto["forma_pago"] = "cuenta corriente"
            gastos.append(gasto)
        elif din_out:
            # Gasto operativo: la categoría se infiere del detalle (los alias
            # del catálogo matchean "luz", "alquiler", etc.; sin match → OTHER
            # con el detalle preservado como label).
            gastos.append({**base, "monto": str(din_out), "categoria": detalle})
        else:
            # Salida de mercadería sin plata (consumo propio / merma / ajuste):
            # no es venta ni gasto → revisión manual.
            otros.append({**base, "mercaderia_egreso": str(mer_out)})

    return {"ventas": ventas, "gastos": gastos, "otros": otros}


def _append_libro_diario_contexts(
    summary: dict[str, Any],
    contexts: list[dict[str, Any]],
    sheet_label: str,
    parsed: dict[str, list[dict[str, Any]]],
    ventas_bucket: list[dict[str, Any]],
    gastos_bucket: list[dict[str, Any]],
) -> None:
    """Vuelca el resultado del Libro Diario en buckets + mapping_contexts.

    Una misma hoja emite hasta tres contextos: ventas, gastos y (si hubo filas
    ambiguas) un contexto sin clasificar que termina en la bandeja "Otros".
    """
    spec = (
        ("ventas", "sale", "Ventas (Dinero/Ingreso)", ventas_bucket,
         ["fecha", "detalle", "monto"]),
        ("gastos", "expense", "Gastos y compras (Dinero/Egreso)", gastos_bucket,
         ["fecha", "detalle", "monto", "categoria", "forma_pago"]),
    )
    for key, entity, sub_label, bucket, headers in spec:
        rows = parsed.get(key) or []
        if not rows:
            continue
        ctx_id = f"sheet:{sheet_label}:{key}"
        bucket.extend({**r, "__context__": ctx_id} for r in rows)
        contexts.append(
            {
                "context_id": ctx_id,
                "label": f"{sheet_label} — {sub_label}",
                "source_kind": "sheet_group",
                "entity_type": entity,
                "headers": headers,
                "fields": None,
                "preview_rows": rows[:10],
                "row_count": len(rows),
                "libro_diario": True,
            }
        )
    otros = parsed.get("otros") or []
    if otros:
        ctx_id = f"sheet:{sheet_label}:otros"
        contexts.append(
            {
                "context_id": ctx_id,
                "label": f"{sheet_label} — Movimientos ambiguos",
                "source_kind": "sheet_group",
                "entity_type": None,
                "headers": None,
                "fields": sorted({k for r in otros for k in r}),
                "preview_rows": otros[:10],
                "row_count": len(otros),
                "unclassified": True,
                "libro_diario": True,
            }
        )
        summary.setdefault("otros_detectados", []).extend(
            {**r, "__context__": ctx_id} for r in otros
        )
        summary["warnings"].append(
            f"{len(otros)} movimiento(s) de '{sheet_label}' son ambiguos (plata que "
            "entra y sale en la misma fila, o salida de mercadería sin plata) y "
            "quedaron para revisión manual."
        )


def _finish_libro_diario_summary(
    summary: dict[str, Any],
    rows: list[list[Any]],
    ld: tuple[int, dict[str, int]],
    label: str,
) -> dict[str, Any]:
    """Completa el summary de un archivo de UNA tabla con formato Libro Diario.

    ``inferred_type="mixed"`` rutea la inserción al path multi-contexto, que
    procesa ventas y gastos por separado desde la misma hoja.
    """
    hdr_idx, col_map = ld
    parsed = parse_libro_diario(rows, hdr_idx, col_map)
    contexts: list[dict[str, Any]] = []
    all_ventas: list[dict[str, Any]] = []
    all_gastos: list[dict[str, Any]] = []
    _append_libro_diario_contexts(summary, contexts, label, parsed, all_ventas, all_gastos)
    total = sum(len(v) for v in parsed.values())
    preview = (parsed["ventas"][:5] + parsed["gastos"][:5])[:10]
    summary.update(
        {
            "libro_diario": True,
            "inferred_type": "mixed",
            "confidence": "HIGH" if (all_ventas or all_gastos) else "LOW",
            "has_venta": bool(all_ventas),
            "has_gasto": bool(all_gastos),
            "has_producto": False,
            "headers": ["fecha", "detalle", "monto"],
            "columns": ["fecha", "detalle", "monto"],
            "row_count": total,
            "rows_processed": total,
            "ventas_detectadas": all_ventas,
            "gastos_detectados": all_gastos,
            "stock_detectado": [],
            "preview_rows": preview,
            "mapping_contexts": contexts,
        }
    )
    return summary


def _decode_text_bytes(content: bytes) -> str:
    """Decodifica bytes de CSV/texto probando encodings comunes (UTF-8 BOM, latin-1, cp1252)."""
    for enc in ("utf-8-sig", "utf-8", "cp1252", "latin-1"):
        try:
            return content.decode(enc)
        except UnicodeDecodeError:
            continue
    return content.decode("utf-8", errors="replace")


def _sniff_delimiter(sample: str) -> str:
    """Detecta el delimitador de un CSV (`,` `;` tab `|`). Default `,` si no puede."""
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        return dialect.delimiter
    except csv.Error:
        # Heurística simple: el candidato más frecuente en la primera línea.
        first_line = sample.splitlines()[0] if sample.splitlines() else ""
        counts = {d: first_line.count(d) for d in (";", ",", "\t", "|")}
        best = max(counts, key=lambda d: counts[d])
        return best if counts[best] > 0 else ","


def classify_line(line: str) -> str:
    # Sin tildes: una línea suelta de un .txt escribe "mercadería" tanto como
    # "mercaderia", y los *_CTX se declaran en la forma sin acento.
    low = strip_accents(line.lower())
    if any(k in low for k in VENTA_CTX):
        return "venta"
    if any(k in low for k in GASTO_CTX):
        return "gasto"
    if any(k in low for k in STOCK_CTX):
        return "stock"
    return "desconocido"


def extract_amounts_from_text(text: str) -> dict[str, Any]:
    lines = text.splitlines()
    ventas: list[dict[str, Any]] = []
    gastos: list[dict[str, Any]] = []
    stock: list[dict[str, Any]] = []

    for line in lines:
        matches = AMOUNT_RE.findall(line)
        if not matches:
            continue
        category = classify_line(line)
        entry = {"linea": line.strip(), "montos": matches}
        if category == "venta":
            ventas.append(entry)
        elif category == "gasto":
            gastos.append(entry)
        elif category == "stock":
            stock.append(entry)
        else:
            ventas.append(entry)

    return {
        "ventas_detectadas": ventas,
        "gastos_detectados": gastos,
        "stock_detectado": stock,
    }


def _store_rows_by_type(
    summary: dict[str, Any],
    rows: list[dict[str, Any]],
    inferred_type: str,
) -> None:
    """Almacena las filas en la clave correcta según el tipo inferido.

    Cada bucket recibe filas SOLO si corresponde a su tipo. `insert_confirmed_data`
    ya cae a `gastos_detectados` cuando `ventas_detectadas` está vacío, así que NO
    hace falta contaminar `ventas_detectadas` con filas de gastos (eso causaba
    riesgo de doble conteo venta+gasto).

    FASE F: el tipo ambiguo ("general") va a `otros_detectados` — nunca más datos
    ambiguos al bucket de ventas por default. El usuario los reasigna al confirmar
    (la inserción legacy también lee este bucket) o quedan en la bandeja "Otros"
    (`unclassified_records`).
    """
    if inferred_type == "stock":
        summary["stock_detectado"] = rows
        summary["ventas_detectadas"] = []
        summary["gastos_detectados"] = []
    elif inferred_type == "gastos":
        summary["gastos_detectados"] = rows
        summary["ventas_detectadas"] = []
        summary["stock_detectado"] = []
    elif inferred_type == "ventas":
        summary["ventas_detectadas"] = rows
        summary["gastos_detectados"] = []
        summary["stock_detectado"] = []
    elif inferred_type == "clientes":
        # F7a: aditivo — el import/vinculación de clientes queda para 7b/7c.
        summary["clientes_detectados"] = rows
        summary["ventas_detectadas"] = []
        summary["gastos_detectados"] = []
        summary["stock_detectado"] = []
    elif inferred_type == "proveedores":
        summary["proveedores_detectados"] = rows
        summary["ventas_detectadas"] = []
        summary["gastos_detectados"] = []
        summary["stock_detectado"] = []
    else:
        summary["otros_detectados"] = rows
        summary["ventas_detectadas"] = []
        summary["gastos_detectados"] = []
        summary["stock_detectado"] = []


# Mapeo tipo-inferido / clasificación-de-hoja → entity_type del ColumnMapper.
_TYPE_TO_ENTITY: dict[str, str] = {
    "ventas": "sale",
    "gastos": "expense",
    "stock": "product",
    "clientes": "customer",
    "proveedores": "supplier",
}


def _build_table_context(
    inferred_type: str,
    headers: list[str],
    preview_rows: list[dict[str, Any]],
    row_count: int,
) -> dict[str, Any]:
    """Construye un único mapping_context para archivos de una tabla (csv / xlsx 1 hoja)."""
    return {
        "context_id": "table",
        "label": "Tabla",
        "source_kind": "table",
        "entity_type": _TYPE_TO_ENTITY.get(inferred_type),
        "headers": list(headers),
        "fields": None,
        "preview_rows": preview_rows,
        "row_count": row_count,
    }


def _build_text_contexts(
    detected: dict[str, list[dict[str, Any]]], source_kind: str
) -> list[dict[str, Any]]:
    """Construye mapping_contexts por grupo detectado en documentos de texto/imagen.

    Estos documentos NO tienen columnas: el "mapeo" es asignar cada grupo de líneas
    detectadas (ventas/gastos/stock) a un entity_type. Marca cada fila con
    `__context__` para separarla en la inserción. `headers=None` señala al frontend
    que muestre preview de líneas + selector de tipo, sin dropdowns de columna.
    """
    # Solo ventas y gastos: una línea de texto (`{linea, montos}`) no alcanza para
    # formar un producto, así que NO se emite contexto para stock_detectado. Esas
    # filas quedan en el summary pero no son importables (igual que el path legacy,
    # que nunca insertó stock desde texto).
    groups = [
        ("ventas_detectadas", "sale", "Ventas detectadas"),
        ("gastos_detectados", "expense", "Gastos detectados"),
    ]
    contexts: list[dict[str, Any]] = []
    for key, entity, label in groups:
        rows = detected.get(key) or []
        if not rows:
            continue
        ctx_id = f"text:{entity}"
        for r in rows:
            r["__context__"] = ctx_id
        fields = sorted({k for r in rows for k in r if not k.startswith("__")})
        contexts.append(
            {
                "context_id": ctx_id,
                "label": label,
                "source_kind": source_kind,
                "entity_type": entity,
                "headers": None,
                "fields": fields,
                "preview_rows": [
                    {k: v for k, v in r.items() if not k.startswith("__")} for r in rows[:10]
                ],
                "row_count": len(rows),
            }
        )
    return contexts


def parse_uploaded_content(content: bytes, mime: str, filename: str) -> dict[str, Any]:
    """Parse uploaded file bytes into a summary compatible with chat and ingestion."""
    if mime in SPREADSHEET_MIMES:
        return _parse_spreadsheet(content, mime, filename)
    if mime in IMAGE_MIMES:
        return _parse_image(content, mime, filename)
    return _parse_document(content, mime, filename)


def _parse_spreadsheet(content: bytes, mime: str, filename: str) -> dict[str, Any]:
    source_format = infer_source_format(filename, mime)
    summary: dict[str, Any] = {
        "file_type": "spreadsheet",
        "source_format": source_format,
        "warnings": [],
    }

    if mime == "text/csv":
        # Encoding robusto (UTF-8 BOM, latin-1, cp1252) + delimitador auto (, ; \t |).
        text = _decode_text_bytes(content)
        delimiter = _sniff_delimiter(text[:8192])
        reader = csv.reader(io.StringIO(text), delimiter=delimiter)
        rows = [r for r in reader if any((c or "").strip() for c in r)]  # descarta filas vacías
        if not rows:
            summary.update(
                {
                    "confidence": "MEDIUM",
                    "ventas_detectadas": [],
                    "rows_processed": 0,
                    "row_count": 0,
                    "columns": [],
                    "preview_rows": [],
                }
            )
            return summary

        # Libro Diario (doble encabezado Dinero/Mercadería × Ingreso/Egreso).
        ld = detect_libro_diario_header(rows)
        if ld is not None:
            return _finish_libro_diario_summary(summary, rows, ld, "Libro diario")

        # Detecta la fila de encabezado real (salta títulos/filas en blanco arriba).
        hdr_idx = _detect_header_row(rows)
        headers = rows[hdr_idx]
        data_rows = rows[hdr_idx + 1 :]
        analysis = analyze_headers(headers)
        all_dicts = rows_to_dicts(headers, data_rows)
        preview_rows = all_dicts[:10]
        null_stats = compute_column_null_stats(all_dicts)
        at_risk = flag_columns_at_risk(null_stats)
        summary.update(analysis)
        summary["delimiter"] = delimiter
        summary["headers"] = headers
        summary["rows_processed"] = len(data_rows)
        summary["row_count"] = len(data_rows)
        summary["columns"] = headers
        summary["preview_rows"] = preview_rows
        if at_risk:
            summary["columns_at_risk"] = at_risk
            summary["warnings"].append(
                f"{len(at_risk)} columna(s) con más del 35% de datos vacíos."
            )
        csv_inferred = analysis.get("inferred_type", "general")
        # Sin truncamiento: todas las filas se guardan (el [:50] previo perdía datos).
        _store_rows_by_type(summary, all_dicts, csv_inferred)
        summary["mapping_contexts"] = [
            _build_table_context(csv_inferred, headers, preview_rows, len(data_rows))
        ]
        return summary

    import unicodedata  # noqa: PLC0415

    import openpyxl  # noqa: PLC0415

    # Sin límite de filas: se importan todas las filas del archivo.
    # JSONB de Neon soporta hasta ~255 MB; un archivo típico de un negocio es < 5 MB.
    _max_rows_per_type = None  # None = sin límite

    def _classify_sheet(name: str) -> str:
        """Clasifica una pestaña de xlsx por nombre →
        'ventas'|'gastos'|'stock'|'clientes'|'proveedores'|'unknown'."""
        norm = unicodedata.normalize("NFD", name.lower().strip())
        norm = "".join(c for c in norm if unicodedata.category(c) != "Mn")
        if any(k in norm for k in ["venta", "ingreso", "cobro", "sale"]):
            return "ventas"
        # Compras de mercadería son ingreso de inventario → productos/stock,
        # NO gasto. (Va antes que la regla de gastos para que "compras a
        # proveedores" se rutee a stock por el match de "compra".)
        if any(k in norm for k in ["compra", "mercaderia", "purchase"]):
            return "stock"
        if any(k in norm for k in ["gasto", "egreso", "operativo", "expense"]):
            return "gastos"
        # F7a: maestros de clientes/proveedores. Van DESPUÉS de compra/gasto para
        # no pisar "Compras a Proveedores" (→ stock) ni "Gastos y Proveedores"
        # (→ gastos) — solo captura un nombre de hoja que es puramente el maestro
        # (ej. "Proveedores"). El contenido (infer_spreadsheet_type) manda siempre;
        # esto es solo el desempate por nombre cuando el contenido es ambiguo.
        if any(k in norm for k in ["cliente", "clientes", "consumidor"]):
            return "clientes"
        if any(k in norm for k in ["proveedor", "proveedores"]):
            return "proveedores"
        if any(k in norm for k in ["producto", "inventario", "catalogo", "stock", "item"]):
            return "stock"
        return "unknown"

    workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    try:
        sheet_names = workbook.sheetnames
        is_multisheet = len(sheet_names) > 1

        if is_multisheet:
            # ── Multi-hoja: combinar datos de cada pestaña por tipo ──────────
            all_ventas: list[dict[str, Any]] = []
            all_gastos: list[dict[str, Any]] = []
            all_stock: list[dict[str, Any]] = []
            # F7a: aditivo — el import/vinculación de clientes/proveedores queda para 7b/7c.
            all_clientes: list[dict[str, Any]] = []
            all_proveedores: list[dict[str, Any]] = []
            primary_headers: list[str] = []
            contexts: list[dict[str, Any]] = []
            total_rows = 0

            # Pre-pass: materializar cada hoja junto con su detección de Libro
            # Diario, para no volver a leer el workbook durante la clasificación.
            sheets_data: list[tuple[str, list[list[Any]], tuple[int, dict[str, int]] | None]] = []
            for sheet_name in sheet_names:
                ws = workbook[sheet_name]
                rows = [list(r) for r in ws.iter_rows(values_only=True)]
                if len(rows) < 2:
                    continue
                sheets_data.append((sheet_name, rows, detect_libro_diario_header(rows)))

            for sheet_name, rows, ld in sheets_data:
                if ld is not None:
                    hdr_idx, col_map = ld
                    parsed = parse_libro_diario(rows, hdr_idx, col_map)
                    _append_libro_diario_contexts(
                        summary, contexts, sheet_name, parsed, all_ventas, all_gastos
                    )
                    total_rows += sum(len(v) for v in parsed.values())
                    continue
                if _ld_norm_cell(sheet_name).startswith(_DERIVED_SHEET_PREFIXES):
                    # Hoja derivada (resumen/ganancias/balance): es un agregado de
                    # los movimientos que ya vienen en el archivo, así que importarla
                    # los contaría dos veces. Se preserva sin clasificar por si el
                    # usuario la quiere.
                    _hdr = _detect_header_row(rows)
                    _headers = [
                        str(c) if c is not None else f"col_{i}"
                        for i, c in enumerate(rows[_hdr])
                    ]
                    _dicts = rows_to_dicts(_headers, rows[_hdr + 1 :])
                    context_id = f"sheet:{sheet_name}"
                    contexts.append(
                        {
                            "context_id": context_id,
                            "label": sheet_name,
                            "source_kind": "sheet",
                            "entity_type": None,
                            "headers": _headers,
                            "fields": None,
                            "preview_rows": _dicts[:10],
                            "row_count": len(_dicts),
                            "unclassified": True,
                        }
                    )
                    summary.setdefault("otros_detectados", []).extend(
                        {**d, "__context__": context_id} for d in _dicts
                    )
                    summary["warnings"].append(
                        f"La hoja '{sheet_name}' es un resumen que Véktor calcula solo "
                        "desde tus movimientos. Importarla sumaría esos totales otra "
                        "vez, encima de los movimientos del mismo archivo. Por eso no "
                        "se importa; podés asignarle una sección a mano si de verdad "
                        "la necesitás."
                    )
                    continue

                # Detecta la fila de encabezado real de la hoja.
                hdr_idx = _detect_header_row(rows)
                headers = [
                    str(c) if c is not None else f"col_{i}"
                    for i, c in enumerate(rows[hdr_idx])
                ]
                data_rows = rows[hdr_idx + 1 :]
                if not data_rows:
                    continue

                dicts = rows_to_dicts(headers, data_rows)
                if not primary_headers:
                    primary_headers = headers

                total_rows += len(data_rows)

                # Content-first: el CONTENIDO (columnas) decide el tipo de hoja; el
                # NOMBRE es solo una orientación/contexto que desempata cuando el
                # contenido es ambiguo. La estructura de datos del usuario no debe
                # condicionar la nuestra: una hoja "Proveedores y Stock" que en realidad
                # es un catálogo (Productos/Stock/Precio de compra/Precio de venta) se
                # rutea como producto por sus columnas, NO como gasto por su nombre.
                # `analyze_headers` ya distingue libro de compras→gastos (regla -1 de
                # `infer_spreadsheet_type`), así que las compras de mercadería siguen
                # cayendo bien aunque traigan columnas de catálogo.
                content_type = analyze_headers(headers).get("inferred_type", "general")
                name_hint = _classify_sheet(sheet_name)
                if content_type != "general":
                    sheet_type = content_type  # el CONTENIDO manda
                elif name_hint != "unknown":
                    sheet_type = name_hint  # el nombre solo desempata si el contenido es ambiguo
                else:
                    sheet_type = "general"  # ambiguo → unclassified / Otros (comportamiento actual)

                entity = _TYPE_TO_ENTITY.get(sheet_type)
                context_id = f"sheet:{sheet_name}"
                if entity is None:
                    # Hoja no clasificable: NO se descarta en silencio. Se preserva como
                    # contexto `unclassified` (entity_type=None) con warning, para que el
                    # usuario la vea y la reclasifique manualmente. El insert salta los
                    # contextos sin entidad, así que no se importa hasta que se asigne tipo.
                    contexts.append(
                        {
                            "context_id": context_id,
                            "label": sheet_name,
                            "source_kind": "sheet",
                            "entity_type": None,
                            "headers": headers,
                            "fields": None,
                            "preview_rows": dicts[:10],
                            "row_count": len(dicts),
                            "unclassified": True,
                        }
                    )
                    # FASE F: las filas completas se preservan en otros_detectados
                    # (antes solo quedaban las 10 de preview) → al confirmar, lo no
                    # reclasificado va a la bandeja "Otros" (unclassified_records).
                    summary.setdefault("otros_detectados", []).extend(
                        {**d, "__context__": context_id} for d in dicts
                    )
                    summary["warnings"].append(
                        f"La hoja '{sheet_name}' no se pudo clasificar automáticamente. "
                        "Revisala y asignale un tipo manualmente si querés importarla."
                    )
                    continue

                contexts.append(
                    {
                        "context_id": context_id,
                        "label": sheet_name,
                        "source_kind": "sheet",
                        "entity_type": entity,
                        "headers": headers,
                        "fields": None,
                        "preview_rows": dicts[:10],
                        "row_count": len(dicts),
                    }
                )
                # Marcador de origen: permite separar filas por contexto en la inserción
                # cuando varias hojas comparten tipo. Se ignora en el mapeo por keyword.
                marked = [{**d, "__context__": context_id} for d in dicts]
                if sheet_type == "ventas":
                    all_ventas.extend(marked)
                elif sheet_type == "gastos":
                    all_gastos.extend(marked)
                elif sheet_type == "stock":
                    all_stock.extend(marked)
                elif sheet_type == "clientes":
                    all_clientes.extend(marked)
                elif sheet_type == "proveedores":
                    all_proveedores.extend(marked)

            has_ventas = bool(all_ventas)
            has_gastos = bool(all_gastos)
            has_stock = bool(all_stock)
            has_clientes = bool(all_clientes)
            has_proveedores = bool(all_proveedores)

            if has_ventas and has_gastos and has_stock or has_ventas and has_gastos:
                inferred_type = "mixed"
            elif has_ventas:
                inferred_type = "ventas"
            elif has_gastos:
                inferred_type = "gastos"
            elif has_stock:
                inferred_type = "stock"
            elif has_clientes and has_proveedores:
                inferred_type = "mixed"
            elif has_clientes:
                inferred_type = "clientes"
            elif has_proveedores:
                inferred_type = "proveedores"
            else:
                inferred_type = "general"

            # Preview global desde los contexts (sin el marcador __context__).
            preview: list[dict[str, Any]] = []
            for _ctx in contexts:
                preview.extend(_ctx["preview_rows"])
            preview = preview[:10]

            summary.update(
                {
                    "multi_sheet": True,
                    "sheet_names": sheet_names,
                    "inferred_type": inferred_type,
                    "confidence": "HIGH" if (has_ventas or has_gastos or has_stock) else "LOW",
                    "has_venta": has_ventas,
                    "has_gasto": has_gastos,
                    "has_producto": has_stock,
                    "headers": primary_headers,
                    "columns": primary_headers,
                    "row_count": total_rows,
                    "rows_processed": total_rows,
                    "ventas_detectadas": all_ventas,
                    "gastos_detectados": all_gastos,
                    "stock_detectado": all_stock,
                    "clientes_detectados": all_clientes,
                    "proveedores_detectados": all_proveedores,
                    "preview_rows": preview,
                    "mapping_contexts": contexts,
                }
            )
            return summary

        # ── Una sola hoja: comportamiento original (con límite ampliado) ─────
        worksheet = workbook.active
        all_rows = list(worksheet.iter_rows(values_only=True))
        if not all_rows:
            summary.update(
                {
                    "confidence": "MEDIUM",
                    "ventas_detectadas": [],
                    "rows_processed": 0,
                    "row_count": 0,
                    "columns": [],
                    "preview_rows": [],
                    "inferred_type": "general",
                }
            )
            return summary

        _all = [list(r) for r in all_rows]

        # Libro Diario (doble encabezado Dinero/Mercadería × Ingreso/Egreso).
        ld = detect_libro_diario_header(_all)
        if ld is not None:
            return _finish_libro_diario_summary(
                summary, _all, ld, str(worksheet.title or "Libro diario")
            )

        # Detecta la fila de encabezado real (salta títulos/filas en blanco arriba).
        hdr_idx = _detect_header_row(_all)
        headers = [
            str(cell) if cell is not None else f"col_{index}"
            for index, cell in enumerate(_all[hdr_idx])
        ]
        data_rows = _all[hdr_idx + 1 :]
        analysis = analyze_headers(headers)
        all_dicts = rows_to_dicts(headers, data_rows)
        preview_rows = all_dicts[:10]
        null_stats = compute_column_null_stats(all_dicts)
        at_risk = flag_columns_at_risk(null_stats)
        summary.update(analysis)
        summary["headers"] = headers
        summary["rows_processed"] = len(data_rows)
        summary["row_count"] = len(data_rows)
        summary["columns"] = headers
        summary["preview_rows"] = preview_rows
        if at_risk:
            summary["columns_at_risk"] = at_risk
            summary["warnings"].append(
                f"{len(at_risk)} columna(s) con más del 35% de datos vacíos."
            )
        inferred = analysis.get("inferred_type", "general")
        _store_rows_by_type(summary, all_dicts, inferred)
        summary["mapping_contexts"] = [
            _build_table_context(inferred, headers, preview_rows, len(data_rows))
        ]
        return summary
    finally:
        workbook.close()


def _parse_document(content: bytes, mime: str, filename: str) -> dict[str, Any]:
    source_format = infer_source_format(filename, mime)
    raw_text = extract_document_text(content, mime, filename)
    detected = extract_amounts_from_text(raw_text)
    preview_rows = _preview_from_detected_rows(detected)  # limpio (antes de marcar)
    contexts = _build_text_contexts(detected, "text_group")  # marca filas con __context__
    row_count = len([line for line in raw_text.splitlines() if line.strip()])
    return {
        "file_type": "text",
        "source_format": source_format,
        "confidence": "MEDIUM",
        "row_count": row_count,
        "columns": [],
        "preview_rows": preview_rows,
        "raw_text_preview": raw_text[:1200],
        "warnings": [],
        "mapping_contexts": contexts,
        **detected,
    }


def _parse_image(content: bytes, mime: str, filename: str) -> dict[str, Any]:
    source_format = infer_source_format(filename, mime)
    summary: dict[str, Any] = {
        "file_type": "image",
        "source_format": source_format,
        "confidence": "LOW",
        "warnings": [],
    }

    try:
        import pytesseract  # noqa: PLC0415
        from PIL import Image, UnidentifiedImageError  # noqa: PLC0415

        try:
            image = Image.open(io.BytesIO(content))
            raw_text = pytesseract.image_to_string(image, lang="spa+eng")
        except (UnidentifiedImageError, OSError):
            raw_text = ""
            summary["warnings"].append(
                "No se pudo abrir la imagen (formato no soportado por Pillow)."
            )

        detected = extract_amounts_from_text(raw_text)
        preview_rows = _preview_from_detected_rows(detected)  # limpio (antes de marcar)
        contexts = _build_text_contexts(detected, "ocr_group")  # marca filas con __context__
        summary.update(
            {
                "raw_text_preview": raw_text[:1200],
                "row_count": len([line for line in raw_text.splitlines() if line.strip()]),
                "columns": [],
                "preview_rows": preview_rows,
                "mapping_contexts": contexts,
                **detected,
            }
        )
        return summary
    except ImportError:
        summary["error"] = "OCR no disponible en este entorno"
        summary["row_count"] = 0
        summary["columns"] = []
        summary["preview_rows"] = []
        summary["ventas_detectadas"] = []
        summary["gastos_detectados"] = []
        summary["stock_detectado"] = []
        return summary


def extract_document_text(content: bytes, mime: str, filename: str) -> str:
    source_format = infer_source_format(filename, mime)

    if source_format == "docx":
        import docx  # noqa: PLC0415

        document = docx.Document(io.BytesIO(content))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)

    if source_format == "pdf":
        from pypdf import PdfReader  # noqa: PLC0415

        reader = PdfReader(io.BytesIO(content))
        page_text = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(part.strip() for part in page_text if part and part.strip())

    if source_format == "pptx":
        from pptx import Presentation  # noqa: PLC0415

        presentation = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and text.strip():
                    lines.append(text.strip())
        return "\n".join(lines)

    return content.decode("utf-8", errors="replace")


def preview_value_from_summary(summary: dict[str, Any]) -> str:
    """Return a compact preview string from parsed_summary_json."""
    preview_candidates = [
        summary.get("preview_rows"),
        summary.get("data"),
        summary.get("rows"),
        summary.get("ventas_detectadas"),
        summary.get("products"),
        summary.get("margins"),
        summary.get("totals"),
        summary.get("raw_text_preview"),
    ]
    for candidate in preview_candidates:
        if candidate:
            return str(candidate)[:400]
    return ""


def summary_row_count(summary: dict[str, Any]) -> int | str:
    row_count = summary.get("row_count")
    if isinstance(row_count, int):
        return row_count
    legacy_count = summary.get("rows_processed")
    if isinstance(legacy_count, int):
        return legacy_count
    preview_rows = summary.get("preview_rows")
    if isinstance(preview_rows, list):
        return len(preview_rows)
    legacy_rows = summary.get("ventas_detectadas")
    if isinstance(legacy_rows, list):
        return len(legacy_rows)
    return "?"


def summary_columns(summary: dict[str, Any]) -> list[str]:
    columns = summary.get("columns")
    if isinstance(columns, list):
        return [str(value) for value in columns]
    headers = summary.get("headers")
    if isinstance(headers, list):
        return [str(value) for value in headers]
    return []


def _preview_from_detected_rows(detected: dict[str, Any]) -> list[dict[str, Any]]:
    preview: list[dict[str, Any]] = []
    for key in ("ventas_detectadas", "gastos_detectados", "stock_detectado"):
        rows = detected.get(key)
        if not isinstance(rows, list):
            continue
        for row in rows[:3]:
            if isinstance(row, dict):
                preview.append(row)
        if preview:
            break
    return preview
