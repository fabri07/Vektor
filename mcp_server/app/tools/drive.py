from __future__ import annotations

import base64
import io
import json
import re
from pathlib import Path
from typing import Any

import filetype
import httpx
from fastapi import HTTPException, status
from sqlalchemy.ext.asyncio import AsyncSession

from app.auth.service import get_valid_access_token
from app.security import RequestContext

DRIVE_FILES_URL = "https://www.googleapis.com/drive/v3/files"
DRIVE_UPLOAD_URL = "https://www.googleapis.com/upload/drive/v3/files"

EXPORT_MIME_MAP = {
    "application/vnd.google-apps.document": ("text/plain", "txt"),
    "application/vnd.google-apps.spreadsheet": (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "xlsx",
    ),
    "application/vnd.google-apps.presentation": ("application/pdf", "pdf"),
}
SAFE_EXTENSION_TO_MIME = {
    ".csv": "text/csv",
    ".txt": "text/plain",
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
STOP_WORDS = {
    "google", "drive", "archivo", "archivos", "carpeta", "carpetas", "buscar", "busca",
    "buscá", "leer", "lee", "analizar", "analiza", "revisar", "revisa", "quiero",
    "necesito", "ver", "mostrar", "mostrame", "sobre", "para", "del", "de", "la",
    "las", "los", "mis", "mi", "en",
}


def _normalize_drive_query(query: str | None) -> str:
    cleaned = re.sub(r"[^\w.\- ]+", " ", (query or "").lower())
    tokens = [token for token in cleaned.split() if token and token not in STOP_WORDS]
    return " ".join(tokens[:8])


def build_drive_search_query(query: str | None, folder_id: str | None) -> str:
    base_filters = ["trashed = false"]
    if folder_id:
        base_filters.append(f"'{folder_id}' in parents")

    normalized = _normalize_drive_query(query)
    if not normalized:
        return " and ".join(base_filters)

    parts: list[str] = []
    for token in normalized.split():
        safe = token.replace("'", "\\'")
        parts.append(f"name contains '{safe}'")
        parts.append(f"fullText contains '{safe}'")
    return f"{' and '.join(base_filters)} and ({' or '.join(parts)})"


def detect_mime(content: bytes, filename: str) -> str:
    kind = filetype.guess(content[:2048])
    if kind and kind.mime:
        return str(kind.mime)
    return SAFE_EXTENSION_TO_MIME.get(Path(filename).suffix.lower(), "application/octet-stream")


def extract_text(content: bytes, mime: str, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if mime == "text/plain" or ext == ".txt":
        return content.decode("utf-8", errors="replace")
    if mime == "text/csv" or ext == ".csv":
        return content.decode("utf-8", errors="replace")
    if mime == "application/pdf" or ext == ".pdf":
        from pypdf import PdfReader  # noqa: PLC0415

        reader = PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    docx_mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if mime == docx_mime or ext == ".docx":
        import docx  # noqa: PLC0415

        document = docx.Document(io.BytesIO(content))
        return "\n".join(paragraph.text for paragraph in document.paragraphs)
    pptx_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if mime == pptx_mime or ext == ".pptx":
        from pptx import Presentation  # noqa: PLC0415

        presentation = Presentation(io.BytesIO(content))
        texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    texts.append(shape.text)
        return "\n".join(texts)
    xlsx_mime = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if mime == xlsx_mime or ext == ".xlsx":
        import openpyxl  # noqa: PLC0415

        workbook = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        try:
            rows: list[str] = []
            for sheet in workbook.worksheets[:3]:
                rows.append(f"[Hoja] {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = [str(cell) for cell in row if cell is not None]
                    if values:
                        rows.append(", ".join(values))
                    if len(rows) >= 80:
                        break
            return "\n".join(rows)
        finally:
            workbook.close()
    return content.decode("utf-8", errors="replace")


def summarize_content(text: str) -> dict[str, Any]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return {
        "line_count": len(lines),
        "content_preview": "\n".join(lines[:12])[:1800],
        "keywords": list(re.findall(r"\b[\w.-]{4,}\b", text.lower())[:20]),
    }


async def list_files(
    *,
    session: AsyncSession,
    ctx: RequestContext,
    folder_id: str | None,
    query: str | None,
    max_results: int,
) -> dict[str, Any]:
    access_token, _ = await get_valid_access_token(session, ctx)
    params: dict[str, str | int] = {
        "pageSize": max_results,
        "q": build_drive_search_query(query, folder_id),
        "fields": "files(id,name,mimeType,modifiedTime,webViewLink,size,parents),nextPageToken",
        "includeItemsFromAllDrives": "true",
        "supportsAllDrives": "true",
        "orderBy": "modifiedTime desc",
    }
    async with httpx.AsyncClient(timeout=20.0) as client:
        resp = await client.get(
            DRIVE_FILES_URL,
            params=params,
            headers={"Authorization": f"Bearer {access_token}"},
        )
    if resp.status_code == status.HTTP_401_UNAUTHORIZED:
        raise HTTPException(status_code=401, detail="refresh_failed")
    if resp.status_code >= 400:
        raise HTTPException(status_code=400, detail="drive_list_failed")
    payload = resp.json()
    files = [
        {
            "id": item.get("id"),
            "name": item.get("name"),
            "mime_type": item.get("mimeType"),
            "modified_time": item.get("modifiedTime"),
            "web_view_link": item.get("webViewLink"),
            "size": item.get("size"),
            "parents": item.get("parents", []),
        }
        for item in payload.get("files", [])
    ]
    return {"files": files}


async def _get_file_metadata(*, access_token: str, file_id: str) -> dict[str, Any]:
    async with httpx.AsyncClient(timeout=20.0) as client:
        resp = await client.get(
            f"{DRIVE_FILES_URL}/{file_id}",
            params={
                "fields": "id,name,mimeType,modifiedTime,webViewLink,size",
                "supportsAllDrives": "true",
            },
            headers={"Authorization": f"Bearer {access_token}"},
        )
    if resp.status_code == 404:
        raise HTTPException(status_code=404, detail="file_not_found")
    if resp.status_code >= 400:
        raise HTTPException(status_code=400, detail="drive_metadata_failed")
    metadata: dict[str, Any] = resp.json()
    return metadata


async def read_file(
    *,
    session: AsyncSession,
    ctx: RequestContext,
    file_id: str,
) -> dict[str, Any]:
    access_token, _ = await get_valid_access_token(session, ctx)
    metadata = await _get_file_metadata(access_token=access_token, file_id=file_id)
    source_name = str(metadata.get("name", "drive-file"))
    source_mime = str(metadata.get("mimeType", "application/octet-stream"))

    async with httpx.AsyncClient(timeout=30.0) as client:
        if source_mime in EXPORT_MIME_MAP:
            export_mime, target_ext = EXPORT_MIME_MAP[source_mime]
            resp = await client.get(
                f"{DRIVE_FILES_URL}/{file_id}/export",
                params={"mimeType": export_mime},
                headers={"Authorization": f"Bearer {access_token}"},
            )
            download_name = f"{source_name}.{target_ext}" if "." not in source_name else source_name
            parse_mime = export_mime
            transfer_mode = "export"
        else:
            resp = await client.get(
                f"{DRIVE_FILES_URL}/{file_id}",
                params={"alt": "media", "supportsAllDrives": "true"},
                headers={"Authorization": f"Bearer {access_token}"},
            )
            download_name = source_name
            parse_mime = detect_mime(resp.content, download_name)
            transfer_mode = "download"

    if resp.status_code == 403:
        raise HTTPException(status_code=403, detail="insufficient_scope")
    if resp.status_code >= 400:
        raise HTTPException(status_code=400, detail="drive_read_failed")

    content = resp.content
    text = extract_text(content, parse_mime, download_name)
    summary = summarize_content(text)
    return {
        "file_id": file_id,
        "name": source_name,
        "mime_type": source_mime,
        "parsed_mime_type": parse_mime,
        "transfer_mode": transfer_mode,
        "content_preview": summary["content_preview"],
        "line_count": summary["line_count"],
        "keywords": summary["keywords"],
        "web_view_link": metadata.get("webViewLink"),
    }


async def upload_file(
    *,
    session: AsyncSession,
    ctx: RequestContext,
    name: str,
    content_base64: str,
    mime_type: str,
    folder_id: str | None,
) -> dict[str, Any]:
    access_token, _ = await get_valid_access_token(session, ctx)
    try:
        content = base64.b64decode(content_base64)
    except Exception as exc:
        raise HTTPException(status_code=400, detail="validation_error") from exc

    boundary = "vektor-google-drive-upload"
    metadata: dict[str, Any] = {"name": name}
    if folder_id:
        metadata["parents"] = [folder_id]

    body = (
        f"--{boundary}\r\n"
        "Content-Type: application/json; charset=UTF-8\r\n\r\n"
        f"{json.dumps(metadata)}\r\n"
        f"--{boundary}\r\n"
        f"Content-Type: {mime_type}\r\n\r\n"
    ).encode() + content + f"\r\n--{boundary}--".encode()

    async with httpx.AsyncClient(timeout=30.0) as client:
        resp = await client.post(
            f"{DRIVE_UPLOAD_URL}?uploadType=multipart&supportsAllDrives=true",
            content=body,
            headers={
                "Authorization": f"Bearer {access_token}",
                "Content-Type": f"multipart/related; boundary={boundary}",
            },
        )
    if resp.status_code >= 400:
        raise HTTPException(status_code=400, detail="drive_upload_failed")
    payload = resp.json()
    return {
        "id": payload.get("id"),
        "name": payload.get("name"),
        "mime_type": payload.get("mimeType"),
        "web_view_link": payload.get("webViewLink"),
    }
